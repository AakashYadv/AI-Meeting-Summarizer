import os
import traceback
from io import BytesIO
from flask import Flask, request, jsonify, render_template
from dotenv import load_dotenv
from utils.emailer import send_email
from groq import Groq
from docx import Document
import PyPDF2
from pptx import Presentation

load_dotenv()

app = Flask(__name__)

def get_groq_client():
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        raise RuntimeError("GROQ_API_KEY is not set. Add it to your environment.")
    return Groq(api_key=api_key)   # ✅ fixed

# Default model
GROQ_MODEL = os.getenv("GROQ_MODEL", "llama-3.3-70b-versatile")

@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")

# -------------------- SUMMARIZATION --------------------
@app.route("/api/summarize", methods=["POST"])
def summarize():
    data = request.get_json(force=True, silent=True) or {}
    transcript = (data.get("transcript") or "").strip()
    instruction = (
        data.get("instruction")
        or "Summarize succinctly with bullets and action items."
    ).strip()

    if not transcript:
        return jsonify({"error": "Transcript is required"}), 400

    system_prompt = (
        "You are an assistant that turns long meeting transcripts into clean, structured, editable summaries.\n"
        "Follow the user's custom instruction exactly. Always include (when applicable):\n"
        "- Title\n"
        "- Executive Summary (3–6 bullets)\n"
        "- Decisions Made\n"
        "- Action Items (owner, due date if mentioned)\n"
        "- Risks/Blockers\n"
        "- Next Steps\n"
        "Be concise. Use simple Markdown. If dates/owners are missing, leave placeholders."
    )

    client = get_groq_client()
    try:
        completion = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {
                    "role": "user",
                    "content": f"Instruction:\n{instruction}\n\nTranscript:\n{transcript}",
                },
            ],
            temperature=0.2,
            max_tokens=1200,
        )
        content = completion.choices[0].message.content
        return jsonify({"summary_markdown": content})
    except Exception as e:
        error_details = traceback.format_exc()
        print("---- Summarize ERROR ----")
        print(error_details)
        return jsonify({"error": str(e), "trace": error_details}), 500

# -------------------- EMAIL --------------------
@app.route("/api/send-email", methods=["POST"])
def email_summary():
    data = request.get_json(force=True, silent=True) or {}
    recipients = data.get("recipients") or []
    subject = (data.get("subject") or "Meeting Summary").strip()
    body_md = (data.get("body_markdown") or "").strip()

    if not recipients or not isinstance(recipients, list):
        return jsonify({"error": "Recipients must be a non-empty list"}), 400
    if not body_md:
        return jsonify({"error": "Body (Markdown) is required"}), 400

    try:
        send_email(recipients, subject, body_md)
        return jsonify({"ok": True})
    except Exception as e:
        error_details = traceback.format_exc()
        print("---- Email ERROR ----")
        print(error_details)
        return jsonify({"error": str(e), "trace": error_details}), 500

# -------------------- FILE UPLOAD --------------------
ALLOWED_EXT = {"txt", "docx", "pdf", "pptx"}

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXT

@app.route("/api/upload-file", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "Empty filename"}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": "Only .txt, .docx, .pdf, .pptx files are supported"}), 400

    ext = file.filename.rsplit(".", 1)[1].lower()
    text = ""

    try:
        if ext == "txt":
            text = file.read().decode("utf-8", errors="ignore")

        elif ext == "docx":
            doc = Document(BytesIO(file.read()))
            text = "\n".join(p.text for p in doc.paragraphs)

        elif ext == "pdf":
            reader = PyPDF2.PdfReader(BytesIO(file.read()))
            pages = [page.extract_text() or "" for page in reader.pages]
            text = "\n".join(pages)

        elif ext == "pptx":
            prs = Presentation(BytesIO(file.read()))
            slides = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slides.append(shape.text)
            text = "\n".join(slides)

        return jsonify({"text": text.strip()})

    except Exception as e:
        error_details = traceback.format_exc()
        print("---- File Upload ERROR ----")
        print(error_details)
        return jsonify({"error": f"Failed to parse file: {str(e)}"}), 500

# -------------------- GLOBAL ERROR HANDLER --------------------
@app.errorhandler(Exception)
def handle_exception(e):
    return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    port = int(os.getenv("PORT", "8080"))
    app.run(host="0.0.0.0", port=port, debug=True)
